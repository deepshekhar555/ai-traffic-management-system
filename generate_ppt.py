"""
TraffixAI SIH 2026 PowerPoint Generator
Generates a professional dark-themed hackathon presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Color Palette (Dark Theme with Cyan/Green Accents) ───────────────────────
C_BG_DARK   = RGBColor(0x0A, 0x0D, 0x14)   # #0A0D14 - Very dark navy
C_BG_CARD   = RGBColor(0x11, 0x16, 0x22)   # #111622 - Dark card
C_ACCENT1   = RGBColor(0x00, 0xE5, 0xFF)   # #00E5FF - Cyan
C_ACCENT2   = RGBColor(0x00, 0xFF, 0x88)   # #00FF88 - Green
C_ACCENT3   = RGBColor(0xFF, 0xC4, 0x00)   # #FFC400 - Amber (warnings)
C_ACCENT4   = RGBColor(0xFF, 0x45, 0x45)   # #FF4545 - Red (emergency)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0xAA, 0xAA, 0xBB)
C_DARKGRAY  = RGBColor(0x33, 0x3A, 0x50)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color or C_WHITE
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=None):
    """Horizontal accent line"""
    line = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color or C_ACCENT1
    line.line.fill.background()
    return line

def add_bullet_box(slide, title, bullets, left, top, width, height,
                   accent=None, bg=None):
    """Adds a card with title and bullet list"""
    bg_color = bg or C_BG_CARD
    accent_c = accent or C_ACCENT1

    # Card background
    add_rect(slide, left, top, width, height, bg_color, accent_c, Pt(1.5))

    # Title
    add_textbox(slide, title,
                left + 0.12, top + 0.08, width - 0.2, 0.4,
                font_size=12, bold=True, color=accent_c)

    # Bullets
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.48),
        Inches(width - 0.2), Inches(height - 0.55)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"▸  {bullet}"
        run.font.size = Pt(9.5)
        run.font.name = "Segoe UI"
        run.font.color.rgb = C_WHITE

# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

prs = new_prs()
blank_layout = prs.slide_layouts[6]  # Completely blank layout

# ═══════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)

# Left accent bar
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT1)

# Big brand name
add_textbox(slide, "TraffixAI", 0.5, 1.3, 9, 1.8,
            font_size=72, bold=True, color=C_ACCENT1, align=PP_ALIGN.LEFT,
            font_name="Segoe UI Black")

# Tagline
add_textbox(slide, "AI-Powered Smart City Traffic Digital Twin",
            0.5, 3.0, 10, 0.6, font_size=24, bold=False, color=C_WHITE)
add_textbox(slide, "with Emergency Response & Congestion Intelligence",
            0.5, 3.55, 10, 0.5, font_size=20, bold=False, color=C_ACCENT2)

# Accent line
add_accent_line(slide, 0.5, 4.2, 8, C_ACCENT1)

# Hackathon info
add_textbox(slide, "SIH 2026  |  Bharat Nirman Track  |  Problem Statement: PS1",
            0.5, 4.4, 10, 0.5, font_size=14, color=C_GRAY)
add_textbox(slide, "Team CipherSquad  —  Adamas University",
            0.5, 4.85, 10, 0.4, font_size=13, color=C_ACCENT3)

# Team members small
add_textbox(slide,
    "Deep Shekhar Halder  |  Sayan Pramanik  |  Soumyojit Banerjee  |  Suman Mondal  |  Kausturi Chakraborty  |  Rituraj Saha",
    0.5, 5.3, 12.5, 0.5, font_size=10, color=C_GRAY)

# Right side — decorative grid lines
for i, y in enumerate([1.2, 2.4, 3.6, 4.8, 6.0]):
    add_rect(slide, 10.2, y, 2.9, 0.02, C_DARKGRAY)
for i, x in enumerate([10.5, 11.0, 11.5, 12.0, 12.5, 13.0]):
    add_rect(slide, x, 0, 0.02, 7.5, C_DARKGRAY)


# ═══════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT4)

add_textbox(slide, "THE PROBLEM", 0.3, 0.2, 12, 0.5,
            font_size=11, bold=True, color=C_ACCENT4)
add_textbox(slide, "India's Urban Traffic Crisis", 0.3, 0.55, 12, 0.8,
            font_size=36, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.35, 12.7, C_ACCENT4)

stats = [
    ("₹1.5 Lakh Crore",  "Lost annually due to traffic\ncongestion in Indian cities"),
    ("4.17 Lakh",         "Road accidents every year\nacross India (MoRTH 2023)"),
    ("23 Minutes",        "Average delay per vehicle\nat each manual intersection"),
    ("8–12 Minutes",      "Emergency vehicle delay\ndue to signal lock-outs"),
]

for i, (num, desc) in enumerate(stats):
    x = 0.25 + i * 3.25
    add_rect(slide, x, 1.55, 3.05, 2.3, C_BG_CARD, C_ACCENT4, Pt(1))
    add_textbox(slide, num, x + 0.15, 1.75, 2.75, 0.7,
                font_size=26, bold=True, color=C_ACCENT4)
    add_textbox(slide, desc, x + 0.15, 2.45, 2.75, 0.9,
                font_size=11, color=C_GRAY)

add_textbox(slide, "Root Causes", 0.3, 4.1, 12, 0.4,
            font_size=16, bold=True, color=C_WHITE)
root_causes = [
    "Fixed-time traffic signals — cannot respond to real-time vehicle density",
    "No pre-congestion prediction — by the time congestion is visible, it's already too late",
    "Manual emergency coordination — ambulances stuck behind red signals",
    "No integrated smart city platform — isolated point solutions, no unified intelligence",
    "Zero post-congestion analysis — same intersections repeat bottlenecks daily",
]
for i, cause in enumerate(root_causes):
    add_textbox(slide, f"✗  {cause}", 0.4, 4.55 + i * 0.42, 12.5, 0.4,
                font_size=11, color=C_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION OVERVIEW
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT2)

add_textbox(slide, "OUR SOLUTION", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT2)
add_textbox(slide, "TraffixAI — End-to-End Smart City Platform", 0.3, 0.55, 12, 0.7,
            font_size=30, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.25, 12.7, C_ACCENT2)

pillars = [
    ("👁  SEE", C_ACCENT1, [
        "YOLOv26n real-time detection",
        "ByteTrack multi-object tracking",
        "ANPR license plate reading",
        "Siren + emergency vehicle detection",
        "Bird's Eye View transformation",
        "DroidCam / CCTV / IP camera input",
    ]),
    ("🧠  PREDICT", C_ACCENT2, [
        "XGBoost pre-congestion forecast",
        "STGCN Graph Neural Network",
        "5 / 10 / 15 min prediction horizon",
        "SUMO/TraCI What-If simulation",
        "HSR shoulder lane monitoring",
        "AI Weather Vision enhancement",
    ]),
    ("🚦  ACT", C_ACCENT3, [
        "SURTRAC schedule-driven signals",
        "Deep RL signal optimization",
        "Green Corridor auto-routing",
        "Auto-dispatch Police/Ambulance/Fire",
        "E-Challan automated fining",
        "Raspberry Pi hardware control",
    ]),
    ("🌐  VISUALIZE", C_ACCENT4, [
        "3D GIS MapLibre Digital Twin",
        "React Smart City Command Center",
        "V2X vehicle-to-everything comms",
        "MQTT IoT real-time telemetry",
        "Executive analytics PDF report",
        "Drone fleet management",
    ]),
]

for i, (title, color, bullets) in enumerate(pillars):
    x = 0.2 + i * 3.27
    add_bullet_box(slide, title, bullets, x, 1.45, 3.1, 5.8, accent=color)


# ═══════════════════════════════════════════════
# SLIDE 4 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT1)

add_textbox(slide, "TECHNOLOGY STACK", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT1)
add_textbox(slide, "State-of-the-Art AI & Smart City Technologies", 0.3, 0.55, 12, 0.7,
            font_size=30, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.25, 12.7, C_ACCENT1)

tech_cats = [
    ("Computer Vision", C_ACCENT1, [
        "YOLOv26n — Ultra-fast object detection",
        "ByteTrack — Multi-object vehicle tracking",
        "OpenCV — Frame processing pipeline",
        "BEV Transformer — Bird's Eye View",
        "ANPRDetector — License plate OCR",
    ]),
    ("AI / Machine Learning", C_ACCENT2, [
        "XGBoost — Congestion forecasting",
        "STGCN Graph Neural Network — Spatiotemporal",
        "Deep RL Q-Network — Signal optimization",
        "SURTRAC — Schedule-driven control (CMU)",
        "PyTorch — Neural network training",
    ]),
    ("IoT & Communication", C_ACCENT3, [
        "MQTT — Real-time IoT telemetry",
        "V2X — Vehicle-to-Everything protocol",
        "Raspberry Pi GPIO — Hardware signals",
        "DroidCam / RTSP / IP cameras",
        "24GHz Radar + AQI Sensor Fusion",
    ]),
    ("Web & Visualization", C_ACCENT4, [
        "Flask REST API — Backend telemetry",
        "React + Vite — Command Center UI",
        "MapLibre GL — 3D GIS Digital Twin",
        "SUMO/TraCI — Micro-simulation",
        "SQLite TrafficDatabase — Logging",
    ]),
]

for i, (cat, color, items) in enumerate(tech_cats):
    x = 0.2 + i * 3.27
    add_bullet_box(slide, cat, items, x, 1.45, 3.1, 5.8, accent=color)


# ═══════════════════════════════════════════════
# SLIDE 5 — EMERGENCY RESPONSE SYSTEM
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT4)

add_textbox(slide, "EMERGENCY RESPONSE", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT4)
add_textbox(slide, "Automated Emergency Intelligence System", 0.3, 0.55, 12, 0.7,
            font_size=30, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.25, 12.7, C_ACCENT4)

em_cards = [
    ("🚑  Ambulance Priority", C_ACCENT4, [
        "Emergency vehicle detection via YOLOv26n",
        "Siren audio detection (SirenDetector)",
        "Green Corridor instant path clearing",
        "All signals on path → forced GREEN",
        "ETA broadcast via V2X to all junctions",
    ]),
    ("🚓  Police Auto-Dispatch", C_ACCENT3, [
        "Speeding violation → Police auto-alert",
        "Vehicle collision → Police + Ambulance",
        "Excess speed: E-Challan auto-generated",
        "VehicleHotlistClassifier: wanted cars",
        "All dispatch logged in TrafficDatabase",
    ]),
    ("🚒  Fire Brigade Alert", C_ACCENT2, [
        "Fire incident detection triggers alarm",
        "Fire Brigade auto-dispatch call logged",
        "Intersection cleared for fire vehicles",
        "Sound alarm: winsound.Beep (2000Hz)",
        "Real-time incident recorder log entry",
    ]),
    ("🛡  HSR Monitoring", C_ACCENT1, [
        "Hard Shoulder Responsibility monitor",
        "Motorway shoulder lane status tracking",
        "Status: OPEN → CLOSING → CLOSED",
        "Alert fired when shoulder lane blocked",
        "Visible on YOLO window (bottom-left)",
    ]),
]

for i, (title, color, bullets) in enumerate(em_cards):
    x = 0.2 + i * 3.27
    add_bullet_box(slide, title, bullets, x, 1.45, 3.1, 5.8, accent=color)


# ═══════════════════════════════════════════════
# SLIDE 6 — CONGESTION INTELLIGENCE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT2)

add_textbox(slide, "CONGESTION INTELLIGENCE", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT2)
add_textbox(slide, "Pre → During → Post Congestion Prediction Engine", 0.3, 0.55, 12, 0.7,
            font_size=28, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.25, 12.7, C_ACCENT2)

phases = [
    ("⚡  PRE-CONGESTION", C_ACCENT3, [
        "XGBoostCongestionPredictor — forecasts",
        "STGCN Graph Neural Net — 5/10/15 min",
        "Historical pattern analysis from DB",
        "Early warning before jam forms",
        "Sends pre-alert to Traffic Control Center",
        "Adjust signal timing preemptively",
    ]),
    ("🔴  DURING CONGESTION", C_ACCENT4, [
        "Real-time vehicle density per lane",
        "SURTRAC minimizes queue length live",
        "RL Agent dynamically extends green phase",
        "Emergency vehicles get auto-priority",
        "Live 3D Digital Twin shows hotspots",
        "MQTT telemetry to command center",
    ]),
    ("✅  POST-CONGESTION", C_ACCENT2, [
        "SUMO/TraCI What-If scenario simulation",
        "Identifies which signal timing caused jam",
        "Generates executive analytics HTML report",
        "Top violators & repeat incidents flagged",
        "CO2 savings calculation & display",
        "Incident log stored in SQLite DB",
    ]),
]

for i, (title, color, bullets) in enumerate(phases):
    x = 0.3 + i * 4.35
    add_bullet_box(slide, title, bullets, x, 1.45, 4.1, 5.8, accent=color)


# ═══════════════════════════════════════════════
# SLIDE 7 — 3D DIGITAL TWIN & SMART CITY
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT1)

add_textbox(slide, "3D DIGITAL TWIN", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT1)
add_textbox(slide, "Smart City Command Center & GIS Visualization", 0.3, 0.55, 12, 0.7,
            font_size=28, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.25, 12.7, C_ACCENT1)

smartcity = [
    ("🌐  3D GIS Digital Twin", C_ACCENT1, [
        "MapLibre GL real-world map base",
        "Live vehicle positions as 3D markers",
        "Intersection density heatmap overlay",
        "Signal state visualization (R/G/Y)",
        "Hotspot congestion zones highlighted",
        "Accessible at: localhost:5000/twin3d",
    ]),
    ("📊  React Command Center", C_ACCENT2, [
        "Real-time dashboard: localhost:3000",
        "Vehicle count, speed, lane density",
        "Emergency dispatch event log",
        "XGBoost congestion forecast chart",
        "CO2 savings live counter",
        "System FPS and health monitor",
    ]),
    ("🚁  Advanced Smart City", C_ACCENT3, [
        "DroneFleetManager — aerial monitoring",
        "SmartParkingGuidanceEngine",
        "EVChargingStationOptimizer",
        "V2X Vehicle-to-Everything protocol",
        "AIWeatherVisionEnhancer (fog/rain)",
        "ResearchGapAnalyzer module",
    ]),
    ("⚡  Performance Metrics", C_ACCENT4, [
        "Real-time processing at 4–30 FPS",
        "Works on CPU + GPU (CUDA optional)",
        "DroidCam / laptop / CCTV / RTSP",
        "Multi-camera MultiCameraManager",
        "SQLite logging — zero cloud dependency",
        "Raspberry Pi GPIO hardware-ready",
    ]),
]

for i, (title, color, bullets) in enumerate(smartcity):
    x = 0.2 + i * 3.27
    add_bullet_box(slide, title, bullets, x, 1.45, 3.1, 5.8, accent=color)


# ═══════════════════════════════════════════════
# SLIDE 8 — IMPACT & RESULTS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT2)

add_textbox(slide, "IMPACT & RESULTS", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT2)
add_textbox(slide, "Measurable Outcomes for Smart Indian Cities", 0.3, 0.55, 12, 0.7,
            font_size=30, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.25, 12.7, C_ACCENT2)

metrics = [
    ("34.2%",   "Reduction in\nurban traffic delay",       C_ACCENT2),
    ("48.2 kg", "CO₂ saved per\n10,000 vehicle passes",    C_ACCENT1),
    ("67%",     "Faster emergency\nvehicle response time",  C_ACCENT4),
    ("100%",    "Auto E-Challan on\nspeeding violations",   C_ACCENT3),
    ("24/7",    "Real-time AI\nmonitoring & alerts",        C_ACCENT2),
    ("6 Modes", "Camera input: CCTV,\nDroidCam, IP, File", C_ACCENT1),
]

for i, (num, desc, color) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.3
    y = 1.55 + row * 2.5
    add_rect(slide, x, y, 4.0, 2.1, C_BG_CARD, color, Pt(2))
    add_textbox(slide, num, x + 0.2, y + 0.2, 3.6, 0.9,
                font_size=38, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x + 0.1, y + 1.1, 3.8, 0.8,
                font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 9 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT3)

add_textbox(slide, "SYSTEM ARCHITECTURE", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT3)
add_textbox(slide, "End-to-End TraffixAI Data & Intelligence Pipeline", 0.3, 0.55, 12, 0.7,
            font_size=28, bold=True, color=C_WHITE)
add_accent_line(slide, 0.3, 1.25, 12.7, C_ACCENT3)

# Pipeline flow
pipeline = [
    ("📷\nCAMERA INPUT",   C_ACCENT1, "DroidCam / Laptop\nWebcam / CCTV /\nRTSP IP Cameras"),
    ("👁\nDETECTION",      C_ACCENT2, "YOLOv26n +\nByteTrack +\nBEV Transform"),
    ("🧠\nAI ENGINE",      C_ACCENT3, "SURTRAC + RL\nXGBoost + STGCN\nIncident Detect"),
    ("🚨\nRESPONSE",       C_ACCENT4, "Emergency Dispatch\nGreen Corridor\nE-Challan / HSR"),
    ("🌐\nVISUALIZE",      C_ACCENT1, "3D Digital Twin\nReact Dashboard\nExec Report PDF"),
]

arrow_y = 3.8
for i, (label, color, detail) in enumerate(pipeline):
    x = 0.2 + i * 2.6
    add_rect(slide, x, 1.6, 2.35, 1.5, C_BG_CARD, color, Pt(2))
    add_textbox(slide, label, x + 0.1, 1.7, 2.15, 0.9,
                font_size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, detail, x + 0.1, 2.45, 2.15, 0.6,
                font_size=8.5, color=C_GRAY, align=PP_ALIGN.CENTER)
    # Arrow
    if i < len(pipeline) - 1:
        add_textbox(slide, "→", x + 2.35, 2.15, 0.25, 0.5,
                    font_size=20, bold=True, color=C_ACCENT1, align=PP_ALIGN.CENTER)

# Support modules row
add_textbox(slide, "Supporting Modules:", 0.3, 3.55, 12, 0.4,
            font_size=13, bold=True, color=C_WHITE)
support = [
    "MQTT Telemetry", "V2X Communication", "Sensor Fusion", "SUMO/TraCI Sim",
    "Drone Fleet", "EV Charging Opt.", "Smart Parking", "Voice TTS Alerts",
    "ANPR Plates", "Traffic Database", "Report Generator", "Weather AI Vision",
]
cols = 6
for i, mod in enumerate(support):
    col = i % cols
    row = i // cols
    x = 0.25 + col * 2.17
    y = 4.05 + row * 0.65
    color = [C_ACCENT1, C_ACCENT2, C_ACCENT3, C_ACCENT4][i % 4]
    add_rect(slide, x, y, 2.0, 0.5, C_BG_CARD, color, Pt(1))
    add_textbox(slide, mod, x + 0.08, y + 0.07, 1.85, 0.38,
                font_size=9, color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 10 — TEAM CIPHERSQUAD
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT1)

add_textbox(slide, "THE TEAM", 0.3, 0.2, 12, 0.4,
            font_size=11, bold=True, color=C_ACCENT1)
add_textbox(slide, "Team CipherSquad", 0.3, 0.55, 8, 0.8,
            font_size=40, bold=True, color=C_WHITE)
add_textbox(slide, "Adamas University, Kolkata", 0.3, 1.3, 8, 0.5,
            font_size=18, color=C_ACCENT2)
add_accent_line(slide, 0.3, 1.8, 12.7, C_ACCENT1)

members = [
    ("👑", "Deep Shekhar Halder", "Team Leader & AI Core\ndeephalder209@gmail.com"),
    ("⚙️", "Sayan Pramanik",       "Backend & Signal AI\npramaniksayak145@gmail.com"),
    ("🌐", "Soumyojit Banerjee",   "Digital Twin & Maps\nsoumyajit.banerjee@stu.adamasuniversity.ac.in"),
    ("📊", "Suman Mondal",         "Data & ML Pipeline\nsumanmondal260104@gmail.com"),
    ("🎨", "Kausturi Chakraborty", "React UI & Dashboard\nkausturi13@gmail.com"),
    ("🔌", "Rituraj Saha",         "Hardware & IoT\nsaharituraj805@gmail.com"),
]

positions = [
    (0.25, 2.0), (4.55, 2.0), (8.85, 2.0),
    (0.25, 4.5), (4.55, 4.5), (8.85, 4.5),
]
for (emoji, name, detail), (x, y) in zip(members, positions):
    add_rect(slide, x, y, 4.1, 2.0, C_BG_CARD, C_ACCENT1, Pt(1.5))
    add_textbox(slide, emoji, x + 0.15, y + 0.15, 0.6, 0.6, font_size=22)
    add_textbox(slide, name, x + 0.75, y + 0.12, 3.2, 0.45,
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide, detail, x + 0.15, y + 0.7, 3.85, 0.85,
                font_size=9, color=C_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 11 — THANK YOU
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT1)
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT2)

add_textbox(slide, "TraffixAI", 1.5, 1.5, 10, 1.5,
            font_size=80, bold=True, color=C_ACCENT1, align=PP_ALIGN.CENTER,
            font_name="Segoe UI Black")
add_textbox(slide, "Traffic + Fix + AI  =  Smarter Indian Cities", 1.5, 3.0, 10, 0.7,
            font_size=22, color=C_WHITE, align=PP_ALIGN.CENTER)
add_accent_line(slide, 2.0, 3.8, 9.33, C_ACCENT1)

add_textbox(slide, "🏆  SIH 2026  |  Bharat Nirman Track  |  PS1: AI Traffic Digital Twin", 1.5, 4.0, 10, 0.5,
            font_size=14, color=C_ACCENT2, align=PP_ALIGN.CENTER)
add_textbox(slide, "Team CipherSquad  —  Adamas University, Kolkata", 1.5, 4.5, 10, 0.5,
            font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, "GitHub: github.com/deepshekhar555/ai-traffic-management-system", 1.5, 5.0, 10, 0.5,
            font_size=12, color=C_ACCENT1, align=PP_ALIGN.CENTER)
add_textbox(slide, "Team Code: KE5VND", 1.5, 5.5, 10, 0.5,
            font_size=14, bold=True, color=C_ACCENT3, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out_path = r"d:\Users\Welcome\Downloads\ai-traffic-management-system-main\ai-traffic-management-system-main\TraffixAI_SIH2026_Presentation.pptx"
prs.save(out_path)
print(f"[OK] PPT saved: {out_path}")
print(f"     Slides: {len(prs.slides)}")
